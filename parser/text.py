"""Parse text frames, paragraphs, and runs."""

import sqlite3

from pptx.util import Emu, Pt

from .color import resolve_color, hex_from_rgb
from .db import insert_row
from .xml_util import get_text_direction, xpath, NSMAP


def _emu_or_none(val) -> int | None:
    """Convert a Length/Emu value to int, or None."""
    if val is None:
        return None
    return int(val)


def _font_size_centipoints(font) -> int | None:
    """Get font size in centi-points (1800 = 18pt). None if inherited."""
    size = font.size
    if size is None:
        return None
    # font.size is a Length in EMU; convert to centi-points
    # 1 pt = 12700 EMU, centi-points = size_emu / 12700 * 100 = size_emu / 127
    return round(int(size) / 127)


def _extract_defRPr(pPr) -> dict:
    """Extract default run properties from paragraph's defRPr element.

    Returns dict with keys: font_name, font_size, font_bold, font_italic,
    font_color. Values are None if not specified.
    """
    defaults = {
        "font_name": None,
        "font_size": None,
        "font_bold": None,
        "font_italic": None,
        "font_color": None,
    }
    if pPr is None:
        return defaults

    defRPr_elems = xpath(pPr, "a:defRPr")
    if not defRPr_elems:
        return defaults
    defRPr = defRPr_elems[0]

    # Font size: sz attr in hundredths of a point → centi-points
    sz = defRPr.get("sz")
    if sz:
        defaults["font_size"] = int(sz)  # already in hundredths of pt = centi-points

    # Bold
    b = defRPr.get("b")
    if b is not None:
        defaults["font_bold"] = b == "1"

    # Italic
    i = defRPr.get("i")
    if i is not None:
        defaults["font_italic"] = i == "1"

    # Font name from latin typeface
    latin = xpath(defRPr, "a:latin")
    if latin:
        tf = latin[0].get("typeface")
        if tf:
            defaults["font_name"] = tf

    # Color from solidFill
    srgb = xpath(defRPr, "a:solidFill/a:srgbClr")
    if srgb:
        val = srgb[0].get("val", "")
        if len(val) == 6:
            defaults["font_color"] = f"#{val.upper()}"
    else:
        # Check for scheme color
        schm = xpath(defRPr, "a:solidFill/a:schemeClr")
        if schm:
            # Store scheme color name for potential resolution
            defaults["_scheme_clr"] = schm[0].get("val")

    return defaults


def _line_spacing_value(para) -> tuple[float | None, str | None]:
    """Extract line spacing as (value, rule)."""
    ls = para.line_spacing
    if ls is None:
        return None, None
    if isinstance(ls, (int, float)):
        return float(ls), "MULTIPLE"
    # It's a Length object (exact spacing)
    return float(int(ls) / 127) / 100, "EXACT"  # convert EMU to points


def _spacing_emu(val) -> int | None:
    """Convert spacing value to EMU."""
    if val is None:
        return None
    return int(val)


def parse_text_frame(shape, shape_id: int, conn: sqlite3.Connection,
                     theme_colors: dict | None = None) -> int | None:
    """Parse a shape's text frame into text_frames/paragraphs/runs tables.
    Returns text_frame_id or None if no text frame.
    """
    if not shape.has_text_frame:
        return None

    tf = shape.text_frame

    # Text direction from XML (not exposed in python-pptx API)
    text_dir = get_text_direction(shape)

    tf_id = insert_row(conn, "text_frames", {
        "shape_id": shape_id,
        "word_wrap": tf.word_wrap if tf.word_wrap is not None else True,
        "auto_size": tf.auto_size.name if tf.auto_size else None,
        "margin_left": _emu_or_none(tf.margin_left),
        "margin_right": _emu_or_none(tf.margin_right),
        "margin_top": _emu_or_none(tf.margin_top),
        "margin_bottom": _emu_or_none(tf.margin_bottom),
        "vertical_anchor": tf.vertical_anchor.name if tf.vertical_anchor else None,
        "text_direction": text_dir,
    })

    for p_idx, para in enumerate(tf.paragraphs):
        _parse_paragraph(para, p_idx, tf_id, conn, theme_colors)

    return tf_id


def _parse_paragraph(para, p_idx: int, tf_id: int, conn: sqlite3.Connection,
                     theme_colors: dict | None = None):
    """Parse a single paragraph."""
    ls_val, ls_rule = _line_spacing_value(para)

    # Bullet properties
    bullet_type = None
    bullet_char = None
    bullet_color = None
    bullet_size_pct = None

    pf = para._pPr  # paragraph properties XML element
    if pf is not None:
        from .xml_util import xpath
        bu_none = xpath(pf, "a:buNone")
        bu_char = xpath(pf, "a:buChar")
        bu_auto = xpath(pf, "a:buAutoNum")
        if bu_none:
            bullet_type = "NONE"
        elif bu_char:
            bullet_type = "CHAR"
            bullet_char = bu_char[0].get("char")
        elif bu_auto:
            bullet_type = "AUTO_NUMBER"

        bu_clr = xpath(pf, "a:buClr/a:srgbClr")
        if bu_clr:
            val = bu_clr[0].get("val", "")
            if len(val) == 6:
                bullet_color = f"#{val.upper()}"

        bu_sz_pct = xpath(pf, "a:buSzPct")
        if bu_sz_pct:
            raw = bu_sz_pct[0].get("val", "")
            if raw:
                bullet_size_pct = int(raw) / 1000  # stored as thousandths of percent

    # Extract default run properties from defRPr
    def_rpr = _extract_defRPr(pf)

    para_id = insert_row(conn, "paragraphs", {
        "text_frame_id": tf_id,
        "paragraph_index": p_idx,
        "alignment": para.alignment.name if para.alignment else None,
        "level": para.level if para.level else 0,
        "space_before": _spacing_emu(para.space_before),
        "space_after": _spacing_emu(para.space_after),
        "line_spacing": ls_val,
        "line_spacing_rule": ls_rule,
        "bullet_type": bullet_type,
        "bullet_char": bullet_char,
        "bullet_color": bullet_color,
        "bullet_size_pct": bullet_size_pct,
        "indent": _emu_or_none(para.indent) if hasattr(para, "indent") else None,
        "margin_left": _emu_or_none(para.margin_left) if hasattr(para, "margin_left") else None,
    })

    # Iterate XML children to handle both <a:r> runs and <a:br/> line breaks
    a_ns = NSMAP["a"]
    p_elem = para._p  # the underlying XML <a:p> element
    r_idx = 0
    run_iter = iter(para.runs)
    for child in p_elem:
        tag = child.tag
        if tag == f"{{{a_ns}}}r":
            # Regular text run
            run = next(run_iter, None)
            if run:
                _parse_run(run, r_idx, para_id, conn, theme_colors, def_rpr)
                r_idx += 1
        elif tag == f"{{{a_ns}}}br":
            # Line break
            insert_row(conn, "runs", {
                "paragraph_id": para_id,
                "run_index": r_idx,
                "text": "",
                "is_line_break": True,
                # Apply defRPr defaults for consistency
                "font_name": def_rpr.get("font_name"),
                "font_size": def_rpr.get("font_size"),
                "font_bold": def_rpr.get("font_bold"),
                "font_italic": def_rpr.get("font_italic"),
                "font_color": def_rpr.get("font_color"),
            })
            r_idx += 1

    # If no children produced runs (empty paragraph with text), add fallback
    if r_idx == 0 and para.text:
        insert_row(conn, "runs", {
            "paragraph_id": para_id,
            "run_index": 0,
            "text": para.text,
            "font_name": def_rpr.get("font_name"),
            "font_size": def_rpr.get("font_size"),
            "font_bold": def_rpr.get("font_bold"),
            "font_italic": def_rpr.get("font_italic"),
            "font_color": def_rpr.get("font_color"),
        })


def _parse_run(run, r_idx: int, para_id: int, conn: sqlite3.Connection,
               theme_colors: dict | None = None, def_rpr: dict | None = None):
    """Parse a single text run, applying defRPr defaults for missing values."""
    font = run.font
    font_color, font_color_theme, brightness = resolve_color(
        font.color, theme_colors
    )
    if def_rpr is None:
        def_rpr = {}

    # Apply defRPr defaults for properties not set on the run
    font_name = font.name or def_rpr.get("font_name")
    font_size = _font_size_centipoints(font)
    if font_size is None:
        font_size = def_rpr.get("font_size")
    font_bold = font.bold
    if font_bold is None:
        font_bold = def_rpr.get("font_bold")
    font_italic = font.italic
    if font_italic is None:
        font_italic = def_rpr.get("font_italic")
    if font_color is None:
        font_color = def_rpr.get("font_color")

    insert_row(conn, "runs", {
        "paragraph_id": para_id,
        "run_index": r_idx,
        "text": run.text,
        "font_name": font_name,
        "font_size": font_size,
        "font_bold": font_bold,
        "font_italic": font_italic,
        "font_underline": font.underline.name if font.underline and font.underline is not True else (
            "SINGLE" if font.underline is True else None
        ),
        "font_strikethrough": None,  # python-pptx doesn't fully expose this
        "font_color": font_color,
        "font_color_theme": font_color_theme,
        "font_color_brightness": brightness,
        "is_line_break": False,
        "is_field": False,
        "field_type": None,
        "hyperlink_url": run.hyperlink.address if run.hyperlink and run.hyperlink.address else None,
    })
