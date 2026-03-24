"""Parse slides: iterate shapes, extract background and notes."""

import json
import sqlite3

from pptx import Presentation
from pptx.enum.dml import MSO_FILL

from .color import hex_from_rgb
from .db import insert_row
from .image import ImageStore
from .shape import parse_shape


def _extract_background(bg) -> dict:
    """Extract background fill properties."""
    result = {"bg_fill_type": None, "bg_fill_color": None, "bg_fill_json": None}
    if bg is None:
        return result
    try:
        fill = bg.fill
        if fill is None or fill.type is None:
            return result
        if fill.type == MSO_FILL.SOLID:
            result["bg_fill_type"] = "solid"
            try:
                result["bg_fill_color"] = hex_from_rgb(fill.fore_color.rgb)
            except Exception:
                pass
        elif fill.type == MSO_FILL.GRADIENT:
            result["bg_fill_type"] = "gradient"
            stops = []
            try:
                for stop in fill.gradient_stops:
                    stops.append({
                        "position": stop.position,
                        "color": hex_from_rgb(stop.color.rgb) if stop.color and stop.color.rgb else None,
                    })
            except Exception:
                pass
            if stops:
                result["bg_fill_json"] = json.dumps({"stops": stops})
    except Exception:
        pass
    return result


def parse_slide_masters(prs: Presentation, pres_id: int,
                        theme_data: dict, conn: sqlite3.Connection) -> dict:
    """
    Parse slide masters and layouts.
    Returns mapping: {slide_layout_obj_id: (layout_db_id, master_db_id, theme_colors)}.
    """
    layout_map = {}

    for m_idx, master in enumerate(prs.slide_masters):
        t_data = theme_data.get(m_idx, {})
        theme_id = t_data.get("theme_id")
        theme_colors = t_data.get("colors", {})

        bg = _extract_background(master.background)
        master_db_id = insert_row(conn, "slide_masters", {
            "presentation_id": pres_id,
            "theme_id": theme_id,
            "name": master.name if hasattr(master, "name") else None,
            **bg,
        })

        for layout in master.slide_layouts:
            layout_bg = _extract_background(layout.background)
            layout_db_id = insert_row(conn, "slide_layouts", {
                "slide_master_id": master_db_id,
                "name": layout.name if hasattr(layout, "name") else None,
                "layout_type": None,
                **layout_bg,
            })
            layout_map[id(layout)] = (layout_db_id, master_db_id, theme_colors)

    return layout_map


def parse_slides(prs: Presentation, pres_id: int, layout_map: dict,
                 conn: sqlite3.Connection, image_store: ImageStore):
    """Parse all slides in the presentation."""
    for slide_num, slide in enumerate(prs.slides, start=1):
        # Resolve layout
        layout_info = layout_map.get(id(slide.slide_layout), (None, None, {}))
        layout_db_id, master_db_id, theme_colors = layout_info

        bg = _extract_background(slide.background)
        slide_db_id = insert_row(conn, "slides", {
            "presentation_id": pres_id,
            "slide_layout_id": layout_db_id,
            "slide_number": slide_num,
            "name": slide.name if hasattr(slide, "name") else None,
            **bg,
        })

        # Parse notes
        if slide.has_notes_slide:
            try:
                notes_tf = slide.notes_slide.notes_text_frame
                notes_text = notes_tf.text if notes_tf else ""
                if notes_text:
                    insert_row(conn, "slide_notes", {
                        "slide_id": slide_db_id,
                        "text": notes_text,
                    })
            except Exception:
                pass

        # Parse shapes
        for z_idx, shape in enumerate(slide.shapes):
            parse_shape(
                shape, conn, image_store, theme_colors,
                slide_id=slide_db_id,
                z_order=z_idx,
            )
