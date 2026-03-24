"""Parse shapes: position, fill, line, shadow, and dispatch by type."""

import json
import sqlite3

from pptx.enum.shapes import MSO_SHAPE_TYPE

from .color import hex_from_rgb, resolve_color
from .db import insert_row
from .image import ImageStore
from .text import parse_text_frame
from .xml_util import get_flip


def _extract_fill(shape) -> dict:
    """Extract fill properties from a shape."""
    result = {"fill_type": None, "fill_color": None, "fill_opacity": None, "fill_json": None}
    try:
        fill = shape.fill
        if fill is None:
            return result
        ft = fill.type
        if ft is None:
            return result

        from pptx.enum.dml import MSO_FILL
        if ft == MSO_FILL.SOLID:
            result["fill_type"] = "solid"
            try:
                rgb = fill.fore_color.rgb
                result["fill_color"] = hex_from_rgb(rgb)
            except Exception:
                pass
        elif ft == MSO_FILL.GRADIENT:
            result["fill_type"] = "gradient"
            stops = []
            try:
                for stop in fill.gradient_stops:
                    stop_data = {
                        "position": stop.position,
                        "color": hex_from_rgb(stop.color.rgb) if stop.color and stop.color.rgb else None,
                    }
                    stops.append(stop_data)
            except Exception:
                pass
            if stops:
                result["fill_json"] = json.dumps({"stops": stops})
        elif ft == MSO_FILL.PATTERNED:
            result["fill_type"] = "pattern"
        elif ft == MSO_FILL.PICTURE:
            result["fill_type"] = "image"
        elif ft == MSO_FILL.BACKGROUND:
            result["fill_type"] = "no_fill"
    except Exception:
        pass
    return result


def _extract_line(shape) -> dict:
    """Extract line/outline properties."""
    result = {"line_color": None, "line_width": None, "line_dash_style": None, "line_opacity": None}
    try:
        line = shape.line
        if line is None:
            return result
        if line.fill and line.fill.type is not None:
            try:
                result["line_color"] = hex_from_rgb(line.color.rgb)
            except Exception:
                pass
        if line.width is not None:
            result["line_width"] = int(line.width)
        if line.dash_style is not None:
            result["line_dash_style"] = line.dash_style.name
    except Exception:
        pass
    return result


def _extract_shadow(shape) -> str | None:
    """Extract shadow properties as JSON string."""
    try:
        shadow = shape.shadow
        if shadow is None:
            return None
        if shadow.inherit:
            return None
        data = {}
        if shadow.style is not None:
            data["style"] = shadow.style.name
        if hasattr(shadow, "blur_radius") and shadow.blur_radius is not None:
            data["blur_radius"] = int(shadow.blur_radius)
        return json.dumps(data) if data else None
    except Exception:
        return None


def _placeholder_info(shape) -> tuple[str | None, int | None]:
    """Extract placeholder type and index."""
    if not shape.is_placeholder:
        return None, None
    ph_format = shape.placeholder_format
    ph_type = ph_format.type.name if ph_format.type else None
    ph_idx = ph_format.idx
    return ph_type, ph_idx


def _shape_type_str(shape) -> str:
    """Map python-pptx shape type to our shape_type string."""
    st = shape.shape_type
    if st == MSO_SHAPE_TYPE.TEXT_BOX:
        return "textbox"
    if st == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if st == MSO_SHAPE_TYPE.TABLE:
        return "table"
    if st == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if st in (MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM):
        # Connectors show up as LINE type
        if hasattr(shape, "begin_x"):
            return "connector"
        if st == MSO_SHAPE_TYPE.FREEFORM:
            return "freeform"
        return "connector"
    if st == MSO_SHAPE_TYPE.CHART:
        return "chart"
    if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return "shape"
    if st == MSO_SHAPE_TYPE.PLACEHOLDER:
        # Determine actual type from placeholder content
        if shape.has_text_frame:
            return "textbox"
        return "shape"
    return "shape"


def _extract_series_color(series) -> str | None:
    """Extract the fill color of a chart series as hex string."""
    try:
        fmt = series.format
        if fmt and fmt.fill and fmt.fill.type is not None:
            rgb = fmt.fill.fore_color.rgb
            return hex_from_rgb(rgb)
    except Exception:
        pass
    return None


def _extract_chart_data(shape) -> str | None:
    """Extract chart data as JSON string from a chart shape."""
    try:
        if not hasattr(shape, "has_chart") or not shape.has_chart:
            return None
        chart = shape.chart
        chart_type_raw = str(chart.chart_type)
        # Normalize chart type name: "COLUMN_CLUSTERED (51)" -> "column_clustered"
        chart_type = chart_type_raw.split("(")[0].strip().lower().replace(" ", "_")

        # Title
        title = None
        try:
            if chart.has_title:
                title = chart.chart_title.text_frame.text
        except Exception:
            pass

        # Legend position
        legend_position = None
        try:
            if chart.has_legend:
                legend_position = chart.legend.position.name.lower()
        except Exception:
            pass

        # Series
        series_list = []
        for series in chart.series:
            s_data = {
                "name": series.name if hasattr(series, "name") else None,
                "values": [v if v is not None else 0 for v in series.values],
            }
            # Extract series color
            color = _extract_series_color(series)
            if color:
                s_data["color"] = color
            # For scatter charts, extract x values from XML
            try:
                ser_el = series._element
                ns = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart"}
                x_val = ser_el.findall(".//c:xVal//c:pt", ns)
                if x_val:
                    x_vals = []
                    for pt in x_val:
                        v_el = pt.find("c:v", ns)
                        x_vals.append(float(v_el.text) if v_el is not None else 0)
                    s_data["x_values"] = x_vals
            except Exception:
                pass
            series_list.append(s_data)

        # Categories
        categories = []
        try:
            plot = chart.plots[0]
            categories = [str(c) for c in plot.categories]
        except Exception:
            pass

        data = {
            "chart_type": chart_type,
            "title": title,
            "categories": categories,
            "series": series_list,
            "legend_position": legend_position,
        }
        return json.dumps(data, ensure_ascii=False)
    except Exception:
        return None


def parse_shape(shape, conn: sqlite3.Connection, image_store: ImageStore,
                theme_colors: dict | None = None,
                slide_id: int | None = None,
                slide_master_id: int | None = None,
                slide_layout_id: int | None = None,
                z_order: int = 0,
                parent_group_id: int | None = None) -> int:
    """Parse a single shape and insert into DB. Returns shape_id."""
    shape_type = _shape_type_str(shape)
    fill = _extract_fill(shape)
    line = _extract_line(shape)
    shadow_json = _extract_shadow(shape)
    ph_type, ph_idx = _placeholder_info(shape)
    flip_h, flip_v = get_flip(shape)

    # Preset geometry
    preset_geo = None
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.auto_shape_type is not None:
            preset_geo = shape.auto_shape_type.name
    except Exception:
        pass

    # Media for pictures
    media_id = None
    crop_left = crop_top = crop_right = crop_bottom = None
    if shape_type == "picture":
        media_id = image_store.get_media_id(shape)
        try:
            crop = shape.crop
            if crop:
                crop_left = int(crop.left * 914400) if crop.left else None
                crop_top = int(crop.top * 914400) if crop.top else None
                crop_right = int(crop.right * 914400) if crop.right else None
                crop_bottom = int(crop.bottom * 914400) if crop.bottom else None
        except Exception:
            pass

    # Connector endpoints
    begin_x = begin_y = end_x = end_y = None
    if shape_type == "connector":
        try:
            begin_x = int(shape.begin_x)
            begin_y = int(shape.begin_y)
            end_x = int(shape.end_x)
            end_y = int(shape.end_y)
        except Exception:
            pass

    # Group child coordinate space
    group_ch_off_x = group_ch_off_y = None
    group_ch_ext_cx = group_ch_ext_cy = None
    if shape_type == "group":
        from .xml_util import get_group_child_transform
        group_ch_off_x, group_ch_off_y, group_ch_ext_cx, group_ch_ext_cy = \
            get_group_child_transform(shape)

    # Chart data
    chart_json = None
    if shape_type == "chart":
        chart_json = _extract_chart_data(shape)

    # Hyperlink
    hyperlink_url = None
    if hasattr(shape, "click_action") and shape.click_action:
        try:
            hyperlink_url = shape.click_action.hyperlink.address
        except Exception:
            pass

    shape_id = insert_row(conn, "shapes", {
        "slide_id": slide_id,
        "slide_master_id": slide_master_id,
        "slide_layout_id": slide_layout_id,
        "shape_type": shape_type,
        "name": shape.name if hasattr(shape, "name") else None,
        "preset_geometry": preset_geo,
        "pos_x": int(shape.left) if shape.left is not None else None,
        "pos_y": int(shape.top) if shape.top is not None else None,
        "width": int(shape.width) if shape.width is not None else None,
        "height": int(shape.height) if shape.height is not None else None,
        "rotation": shape.rotation if hasattr(shape, "rotation") else 0,
        "flip_h": flip_h,
        "flip_v": flip_v,
        "z_order": z_order,
        "parent_group_id": parent_group_id,
        "placeholder_type": ph_type,
        "placeholder_idx": ph_idx,
        **fill,
        **line,
        "shadow_json": shadow_json,
        "media_id": media_id,
        "crop_left": crop_left,
        "crop_top": crop_top,
        "crop_right": crop_right,
        "crop_bottom": crop_bottom,
        "begin_x": begin_x,
        "begin_y": begin_y,
        "end_x": end_x,
        "end_y": end_y,
        "group_ch_off_x": group_ch_off_x,
        "group_ch_off_y": group_ch_off_y,
        "group_ch_ext_cx": group_ch_ext_cx,
        "group_ch_ext_cy": group_ch_ext_cy,
        "hyperlink_url": hyperlink_url,
        "raw_xml_snippet": None,
        "chart_json": chart_json,
    })

    # Parse text frame if present
    if shape_type not in ("picture", "connector", "group", "table"):
        parse_text_frame(shape, shape_id, conn, theme_colors)

    # Parse table if table shape
    if shape_type == "table":
        from .table import parse_table
        parse_table(shape, shape_id, conn, theme_colors)

    # Parse group children
    if shape_type == "group":
        from .group import parse_group_children
        parse_group_children(shape, shape_id, conn, image_store, theme_colors,
                             slide_id=slide_id, slide_master_id=slide_master_id,
                             slide_layout_id=slide_layout_id)

    return shape_id
