"""Parse table shapes: dimensions, cells, and cell text."""

import json
import sqlite3

from .color import hex_from_rgb
from .db import insert_row
from .text import parse_text_frame


def parse_table(shape, shape_id: int, conn: sqlite3.Connection,
                theme_colors: dict | None = None):
    """Parse a table shape into table_dimensions and table_cells."""
    tbl = shape.table
    num_rows = len(tbl.rows)
    num_cols = len(tbl.columns)

    col_widths = [int(col.width) for col in tbl.columns]
    row_heights = [int(row.height) for row in tbl.rows]

    insert_row(conn, "table_dimensions", {
        "shape_id": shape_id,
        "num_rows": num_rows,
        "num_cols": num_cols,
        "col_widths_json": json.dumps(col_widths),
        "row_heights_json": json.dumps(row_heights),
    })

    for r_idx in range(num_rows):
        for c_idx in range(num_cols):
            cell = tbl.cell(r_idx, c_idx)
            _parse_cell(cell, shape_id, r_idx, c_idx, conn, theme_colors)


def _extract_cell_borders(cell) -> dict:
    """Extract border properties from cell XML."""
    borders = {}
    try:
        tc = cell._tc
        for side in ("lnL", "lnR", "lnT", "lnB"):
            side_map = {"lnL": "left", "lnR": "right", "lnT": "top", "lnB": "bottom"}
            db_side = side_map[side]
            from .xml_util import xpath
            ln_elems = xpath(tc, f".//a:tcPr/a:{side}")
            if ln_elems:
                ln = ln_elems[0]
                w = ln.get("w")
                borders[f"border_{db_side}_width"] = int(w) if w else None
                srgb = xpath(ln, ".//a:srgbClr")
                if srgb:
                    val = srgb[0].get("val", "")
                    borders[f"border_{db_side}_color"] = f"#{val.upper()}" if len(val) == 6 else None
                else:
                    borders[f"border_{db_side}_color"] = None
            else:
                borders[f"border_{db_side}_width"] = None
                borders[f"border_{db_side}_color"] = None
    except Exception:
        for side in ("left", "right", "top", "bottom"):
            borders[f"border_{side}_width"] = None
            borders[f"border_{side}_color"] = None
    return borders


def _parse_cell(cell, shape_id: int, r_idx: int, c_idx: int,
                conn: sqlite3.Connection, theme_colors: dict | None = None):
    """Parse a single table cell."""
    # Merge info
    is_merge_origin = cell.is_merge_origin if hasattr(cell, "is_merge_origin") else True
    span_h = cell.span_width if hasattr(cell, "span_width") else 1
    span_v = cell.span_height if hasattr(cell, "span_height") else 1

    # Fill
    fill_type = None
    fill_color = None
    try:
        fill = cell.fill
        if fill and fill.type is not None:
            from pptx.enum.dml import MSO_FILL
            if fill.type == MSO_FILL.SOLID:
                fill_type = "solid"
                fill_color = hex_from_rgb(fill.fore_color.rgb)
    except Exception:
        pass

    # Margins
    margin_left = int(cell.margin_left) if cell.margin_left is not None else None
    margin_right = int(cell.margin_right) if cell.margin_right is not None else None
    margin_top = int(cell.margin_top) if cell.margin_top is not None else None
    margin_bottom = int(cell.margin_bottom) if cell.margin_bottom is not None else None

    # Vertical anchor
    vert_anchor = None
    if cell.vertical_anchor is not None:
        vert_anchor = cell.vertical_anchor.name

    borders = _extract_cell_borders(cell)

    # Parse cell text into text_frames
    # Table cells always have a text_frame (no has_text_frame check needed)
    tf_id = None
    try:
        tf = cell.text_frame

        tf_id = insert_row(conn, "text_frames", {
            "shape_id": shape_id,  # link to the table shape
            "word_wrap": getattr(tf, "word_wrap", True),
            "auto_size": tf.auto_size.name if getattr(tf, "auto_size", None) else None,
            "margin_left": None,
            "margin_right": None,
            "margin_top": None,
            "margin_bottom": None,
            "vertical_anchor": vert_anchor,
            "text_direction": None,
        })

        from .text import _parse_paragraph
        for p_idx, para in enumerate(tf.paragraphs):
            _parse_paragraph(para, p_idx, tf_id, conn, theme_colors)
    except Exception:
        pass

    insert_row(conn, "table_cells", {
        "shape_id": shape_id,
        "row_idx": r_idx,
        "col_idx": c_idx,
        "row_span": span_v,
        "col_span": span_h,
        "is_merge_origin": is_merge_origin,
        "fill_type": fill_type,
        "fill_color": fill_color,
        "margin_left": margin_left,
        "margin_right": margin_right,
        "margin_top": margin_top,
        "margin_bottom": margin_bottom,
        "vertical_anchor": vert_anchor,
        **borders,
        "text_frame_id": tf_id,
    })
